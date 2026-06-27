import os
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb
from chromadb.utils import embedding_functions

CHROMA_DIR = "app/chroma_db"
UPLOAD_DIR = "app/uploads"

chroma_client = chromadb.PersistentClient(path=CHROMA_DIR)
embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name="all-MiniLM-L6-v2"
)


def get_collection():
    return chroma_client.get_or_create_collection(
        name="documents",
        embedding_function=embedding_fn
    )


def extract_text(file_path: str, filename: str) -> str:
    """Extract text from any supported file type."""
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            print(f"Extracted {len(text)} characters from {filename}")
            return text

        elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"]:
            import pytesseract
            from PIL import Image
            pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            print(f"OCR extracted {len(text)} characters from {filename}")
            return text

        elif ext in [".docx", ".doc"]:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
            print(f"Extracted {len(text)} characters from {filename}")
            return text

        elif ext in [".txt", ".csv"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            print(f"Extracted {len(text)} characters from {filename}")
            return text

        elif ext in [".xlsx", ".xls"]:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            rows = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        rows.append(row_text)
            text = "\n".join(rows)
            print(f"Extracted {len(text)} characters from {filename}")
            return text

        elif ext in [".pptx", ".ppt"]:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides.append(shape.text)
            text = "\n".join(slides)
            print(f"Extracted {len(text)} characters from {filename}")
            return text

        else:
            print(f"Unsupported file type: {ext}")
            return ""

    except Exception as e:
        print(f"Text extraction error for {filename}: {e}")
        return ""


def extract_text_from_pdf(file_path: str) -> str:
    """Backward compat alias."""
    return extract_text(file_path, os.path.basename(file_path))


def add_document(file_path: str, filename: str) -> bool:
    try:
        text = extract_text(file_path, filename)
        if not text.strip():
            print(f"No text extracted from {filename}")
            return False

        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        chunks = splitter.split_text(text)
        print(f"Split {filename} into {len(chunks)} chunks")

        col = get_collection()

        # Remove old chunks for this file if re-uploading
        try:
            existing = col.get(where={"filename": filename})
            if existing["ids"]:
                col.delete(ids=existing["ids"])
        except Exception:
            pass

        ids = [f"{filename}_chunk_{i}" for i in range(len(chunks))]
        metadatas = [{"filename": filename, "chunk": i} for i in range(len(chunks))]
        col.add(documents=chunks, ids=ids, metadatas=metadatas)
        print(f"Indexed {len(chunks)} chunks for {filename}")
        return True

    except Exception as e:
        print(f"Error adding document: {e}")
        return False


def add_document_to_vectorstore(file_path: str, filename: str = None) -> list:
    """Alias used by files.py — adds doc and returns stored IDs."""
    fname = filename or os.path.basename(file_path)
    success = add_document(file_path, fname)
    if not success:
        return []
    col = get_collection()
    try:
        existing = col.get(where={"filename": fname})
        return existing["ids"] if existing["ids"] else []
    except Exception:
        return []


def delete_document_from_vectorstore(chroma_ids: list):
    """Delete chunks by ID — used by files.py."""
    if not chroma_ids:
        return
    try:
        col = get_collection()
        col.delete(ids=chroma_ids)
    except Exception as e:
        print(f"Delete error: {e}")


def search_documents(query: str, n_results: int = 5) -> str:
    """Simple search — returns combined text string."""
    try:
        col = get_collection()
        count = col.count()
        print(f"Collection has {count} chunks")
        if count == 0:
            return ""
        results = col.query(
            query_texts=[query],
            n_results=min(n_results, count)
        )
        if not results["documents"] or not results["documents"][0]:
            return ""
        context = "\n\n".join(results["documents"][0])
        print(f"Found {len(results['documents'][0])} relevant chunks")
        return context
    except Exception as e:
        print(f"Search error: {e}")
        return ""


def search_documents_with_scores(query: str, k: int = 4) -> list:
    """Search ChromaDB and return (chunk, distance, filename) tuples.
    Distance: lower = more similar. 0.0 = identical, 2.0 = completely different.
    Only return chunks with distance < 0.75 (relevant threshold).
    """
    try:
        col = get_collection()
        count = col.count()
        if count == 0:
            print("Collection is empty — no documents indexed yet")
            return []

        results = col.query(
            query_texts=[query],
            n_results=min(k, count),
            include=["documents", "distances", "metadatas"]
        )

        docs = results.get("documents", [[]])[0]
        distances = results.get("distances", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]

        if not docs:
            return []

        # Log what was found for debugging
        for doc, dist, meta in zip(docs, distances, metadatas):
            fname = meta.get("filename", "unknown") if meta else "unknown"
            print(f"  Found chunk from '{fname}' — distance: {dist:.3f}")

        return list(zip(docs, distances))

    except Exception as e:
        print(f"Search with scores error: {e}")
        return []


def get_all_documents() -> list:
    try:
        col = get_collection()
        results = col.get()
        if not results["metadatas"]:
            return []
        return list(set(m["filename"] for m in results["metadatas"]))
    except Exception as e:
        print(f"Error getting documents: {e}")
        return []


def delete_document(filename: str) -> bool:
    try:
        col = get_collection()
        existing = col.get(where={"filename": filename})
        if existing["ids"]:
            col.delete(ids=existing["ids"])
            print(f"Deleted {len(existing['ids'])} chunks for {filename}")
            return True
        return False
    except Exception as e:
        print(f"Error deleting: {e}")
        return False
